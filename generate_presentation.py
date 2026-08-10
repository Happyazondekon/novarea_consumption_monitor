from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper function to add a slide with title and bullets
    def add_bullet_slide(title_text, bullets):
        slide_layout = prs.slide_layouts[1] # Bullet layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text

        tf = slide.placeholders[1].text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        # Add a placeholder box for screenshots
        left = Inches(0.5)
        top = Inches(4.5)
        width = Inches(9)
        height = Inches(2.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf_note = txBox.text_frame
        p_note = tf_note.add_paragraph()
        p_note.text = "📷 [ INSERT SCREENSHOT HERE ]"
        p_note.font.bold = True
        p_note.font.size = Pt(24)
        p_note.font.color.rgb = RGBColor(0, 102, 204) # Novarea Blue
        p_note.alignment = PP_ALIGN.CENTER

    # Slide 1: Cover
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Novarea Textiles Monitoring Platform"
    subtitle.text = "Real-Time Consumption Analytics & Industrial Auditing Solution\nOfficial Deployment - August 2026"

    # Slide 2: Objectives
    add_bullet_slide("Streamlining Industrial Efficiency", [
        "Centralized Control: Unified monitoring of Power and Water resources.",
        "Data Integrity: Secure digital audit trail for all consumption logs.",
        "Operational Intelligence: Turning raw indexes into actionable trends.",
        "Accountability: Specialized access for Technicians and Admins."
    ])

    # Slide 3: Roles
    add_bullet_slide("Strategic Access Control", [
        "Field Technicians: Mobile entry, visual evidence, direct missions.",
        "System Administrators: Advanced dashboard, auditing, report generation.",
        "Secure access based on verified factory credentials."
    ])

    # Slide 4: Technician Entry
    add_bullet_slide("Technician Workflow: Data Capture", [
        "3-Step Submission: Resource -> Value -> Confirmation.",
        "Visual Proof: Mandatory photo upload to ensure authenticity.",
        "Instant Sync: Immediate transmission to the Admin dashboard."
    ])

    # Slide 5: Mission Control
    add_bullet_slide("Technician Workflow: Mission Control", [
        "Active Missions: Management instructions received in real-time.",
        "Mark as Done: Automated completion notification to supervisors.",
        "Personal History: Track and verify all past submissions."
    ])

    # Slide 6: Admin Dashboard
    add_bullet_slide("Admin Workflow: The Command Center", [
        "Intra-Day Tracking: Live updates upon second daily reading.",
        "Key Indicators: Last Reading, Usage, Events, and Daily Averages.",
        "Smart Visualization: High-fidelity charts with reference averages."
    ])

    # Slide 7: Linear Interpolation
    add_bullet_slide("Data Integrity: Linear Interpolation", [
        "Gap Management: Automatic logic for missing reading days.",
        "Mathematical Precision: Balanced distribution of measured totals.",
        "Transparency: (*) Indicators for all estimated data points."
    ])

    # Slide 8: Event Logs
    add_bullet_slide("Operational Context: Event Logs", [
        "Contextual Intelligence: Categorizing anomalies (Ops, Mech, Human, Hiding).",
        "Explainable Data: Correlating usage peaks with factory events.",
        "Full Audit Trail: Chronological log of all industrial context."
    ])

    # Slide 9: Report Generation
    add_bullet_slide("Intelligence Unit: Report Generation", [
        "Multi-Format Export: High-fidelity PDF and detailed Excel sheets.",
        "Temporal Flexibility: Weekly, Monthly, Yearly, or Custom ranges.",
        "Modular Sections: Choose between Trends, Logs, or Event summaries."
    ])

    # Slide 10: Conclusion
    add_bullet_slide("Conclusion", [
        "Reliable & Future-Proof Monitoring Architecture.",
        "Empowering data-driven decisions for Novarea Management.",
        "Full Traceability - Excellence through Precision."
    ])

    prs.save('Novarea_Monitoring_Presentation.pptx')
    print("Presentation saved as 'Novarea_Monitoring_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
